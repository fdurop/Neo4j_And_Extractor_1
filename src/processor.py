import os
import json
import zipfile
import tempfile
import shutil
import xml.etree.ElementTree as ET

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import pandas as pd
import re
import subprocess


class AdvancedPPTProcessor:
    def __init__(self, preprocessor, fast_mode=False):
        """
        初始化高级PPTX处理器

        Args:
            preprocessor: MultimodalPreprocessor实例，用于重用输出目录与结果记录
            fast_mode: 快速模式，跳过耗时的CLIP描述生成
        """
        self.preprocessor = preprocessor
        self.fast_mode = fast_mode
        self.output_text_dir = "output/text"
        self.output_table_dir = "output/tables"
        self.output_img_dir = "output/images"

        # 确保输出目录存在
        os.makedirs(self.output_text_dir, exist_ok=True)
        os.makedirs(self.output_table_dir, exist_ok=True)
        os.makedirs(self.output_img_dir, exist_ok=True)

    def extract_all_images_via_zip(self, file_path):
        """
        通过ZIP解压和XML解析提取PPTX中的所有图片

        Args:
            file_path: PPTX文件路径

        Returns:
            dict: 包含幻灯片到图片映射关系的字典
        """
        print(f"开始通过ZIP方式提取图片: {file_path}")

        base_filename = os.path.splitext(os.path.basename(file_path))[0]
        slide_image_mapping = {}

        # 创建临时目录
        with tempfile.TemporaryDirectory() as temp_dir:
            try:
                # 1. 解压PPTX文件
                print("正在解压PPTX文件...")
                with zipfile.ZipFile(file_path, 'r') as zip_ref:
                    zip_ref.extractall(temp_dir)

                # 2. 找到媒体目录
                media_dir = os.path.join(temp_dir, "ppt", "media")
                slides_dir = os.path.join(temp_dir, "ppt", "slides")
                rels_dir = os.path.join(temp_dir, "ppt", "slides", "_rels")

                if not os.path.exists(media_dir):
                    print("未找到media目录，可能没有图片")
                    return slide_image_mapping

                print(f"找到media目录: {media_dir}")
                print(f"媒体文件: {os.listdir(media_dir)}")

                # 3. 遍历所有幻灯片XML文件
                if os.path.exists(slides_dir):
                    for slide_file in os.listdir(slides_dir):
                        if slide_file.startswith("slide") and slide_file.endswith(".xml"):
                            slide_num = self._extract_slide_number(slide_file)
                            if slide_num is None:
                                continue

                            print(f"处理幻灯片 {slide_num}: {slide_file}")

                            # 解析幻灯片XML获取图片关系ID
                            slide_xml_path = os.path.join(slides_dir, slide_file)
                            image_rids = self._parse_slide_xml_for_images(slide_xml_path)

                            if image_rids:
                                print(f"幻灯片 {slide_num} 中找到图片关系ID: {image_rids}")

                                # 解析关系文件获取实际文件名
                                rels_file = slide_file + ".rels"
                                rels_path = os.path.join(rels_dir, rels_file)

                                if os.path.exists(rels_path):
                                    image_files = self._parse_rels_file(rels_path, image_rids)

                                    if image_files:
                                        slide_image_mapping[slide_num] = image_files
                                        print(f"幻灯片 {slide_num} 映射到图片: {image_files}")

                                        # 复制图片到输出目录
                                        self._copy_images_to_output(media_dir, image_files,
                                                                    base_filename, slide_num)

                print(f"图片提取完成，映射关系: {slide_image_mapping}")

            except Exception as e:
                print(f"ZIP方式图片提取失败: {e}")
                import traceback
                traceback.print_exc()

        return slide_image_mapping

    def _extract_slide_number(self, slide_filename):
        """从幻灯片文件名中提取编号"""
        try:
            # slide1.xml -> 1
            import re
            match = re.search(r'slide(\d+)\.xml', slide_filename)
            if match:
                return int(match.group(1))
        except Exception:
            pass
        return None

    def _parse_slide_xml_for_images(self, slide_xml_path):
        """
        解析幻灯片XML文件，查找图片引用

        Args:
            slide_xml_path: 幻灯片XML文件路径

        Returns:
            list: 图片关系ID列表
        """
        image_rids = []

        try:
            tree = ET.parse(slide_xml_path)
            root = tree.getroot()

            # 定义命名空间
            namespaces = {
                'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
                'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
                'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'
            }

            # 查找所有a:blip元素（图片引用）
            blip_elements = root.findall('.//a:blip', namespaces)

            for blip in blip_elements:
                embed_attr = blip.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
                if embed_attr:
                    image_rids.append(embed_attr)
                    print(f"找到图片引用ID: {embed_attr}")

        except Exception as e:
            print(f"解析幻灯片XML失败 {slide_xml_path}: {e}")

        return image_rids

    def _parse_rels_file(self, rels_path, image_rids):
        """
        解析关系文件，获取关系ID到文件名的映射

        Args:
            rels_path: 关系文件路径
            image_rids: 图片关系ID列表

        Returns:
            list: 对应的图片文件名列表
        """
        image_files = []

        try:
            tree = ET.parse(rels_path)
            root = tree.getroot()

            # 定义命名空间
            namespaces = {
                'rel': 'http://schemas.openxmlformats.org/package/2006/relationships'
            }

            # 查找所有关系
            for relationship in root.findall('.//rel:Relationship', namespaces):
                rel_id = relationship.get('Id')
                target = relationship.get('Target')
                rel_type = relationship.get('Type')

                # 检查是否是图片关系
                if (rel_id in image_rids and
                        target and
                        rel_type and
                        'image' in rel_type.lower()):
                    # 提取文件名 (../media/image1.png -> image1.png)
                    filename = os.path.basename(target)
                    image_files.append(filename)
                    print(f"关系映射: {rel_id} -> {filename}")

        except Exception as e:
            print(f"解析关系文件失败 {rels_path}: {e}")

        return image_files

    def _copy_images_to_output(self, media_dir, image_files, base_filename, slide_num):
        """
        将图片复制到输出目录并生成描述

        Args:
            media_dir: 媒体文件源目录
            image_files: 图片文件名列表
            base_filename: 基础文件名
            slide_num: 幻灯片编号
        """
        for idx, image_file in enumerate(image_files, 1):
            try:
                source_path = os.path.join(media_dir, image_file)

                if os.path.exists(source_path):
                    # 生成输出文件名
                    file_ext = os.path.splitext(image_file)[1]
                    output_filename = f"{base_filename}_slide_{slide_num}_img_{idx}_zip{file_ext}"
                    output_path = os.path.join(self.output_img_dir, output_filename)

                    # 复制图片
                    shutil.copy2(source_path, output_path)
                    print(f"复制图片: {source_path} -> {output_path}")

                    # 生成CLIP描述（根据模式决定是否生成）
                    desc_path = None
                    if not self.fast_mode:
                        try:
                            desc_path = self.preprocessor.clip_generate_description(output_path)
                        except Exception as e:
                            print(f"生成CLIP描述失败，跳过: {e}")
                    else:
                        print("快速模式：跳过CLIP描述生成")

                    # 记录到结果中
                    self.preprocessor.results.append({
                        "type": "ppt_image_zip",
                        "page": slide_num,
                        "file": output_path,
                        "description_file": desc_path,
                        "extraction_method": "zip_xml_parsing",
                        "original_filename": image_file
                    })

                else:
                    print(f"源图片文件不存在: {source_path}")

            except Exception as e:
                print(f"复制图片失败 {image_file}: {e}")

    def extract_and_convert_equations(self, slide, slide_number):
        """
        处理幻灯片中的公式

        Args:
            slide: python-pptx的Slide对象
            slide_number: 幻灯片编号

        Returns:
            list: 包含公式信息的列表
        """
        equations = []

        for shape_index, shape in enumerate(slide.shapes):
            try:
                # 检查形状是否包含文本框
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    # 获取形状的XML内容
                    shape_xml = self._get_shape_xml(shape)
                    if shape_xml:
                        # 检查是否包含OMML公式标签
                        omml_content = self._extract_omml_from_xml(shape_xml)
                        if omml_content:
                            print(f"在幻灯片 {slide_number} 形状 {shape_index} 中发现OMML公式")

                            # 尝试转换OMML到LaTeX
                            latex_content = self._convert_omml_to_latex(omml_content)

                            equation_info = {
                                "slide_number": slide_number,
                                "shape_index": shape_index,
                                "type": "omml_formula",
                                "original_omml": omml_content[:500] + "..." if len(
                                    omml_content) > 500 else omml_content,
                                "latex": latex_content,
                                "conversion_success": latex_content is not None
                            }

                            equations.append(equation_info)

                            # 添加到结果中
                            self.preprocessor.results.append({
                                "type": "formula",
                                "page": slide_number,
                                "formula_type": "omml",
                                "latex": latex_content,
                                "source": f"slide_{slide_number}_shape_{shape_index}",
                                "conversion_method": "omml_to_latex"
                            })

                            continue

                # 如果没有找到OMML，检查是否为可能的公式图片
                if self._is_potential_formula_image(shape):
                    print(f"在幻灯片 {slide_number} 形状 {shape_index} 中发现潜在公式图片")

                    # 使用图片处理流程处理公式图片
                    formula_image_path = self._process_formula_image(shape, slide_number, shape_index)

                    if formula_image_path:
                        equation_info = {
                            "slide_number": slide_number,
                            "shape_index": shape_index,
                            "type": "formula_image",
                            "image_path": formula_image_path,
                            "latex": None,
                            "conversion_success": False
                        }

                        equations.append(equation_info)

                        # 添加到结果中
                        self.preprocessor.results.append({
                            "type": "formula",
                            "page": slide_number,
                            "formula_type": "image",
                            "image_path": formula_image_path,
                            "source": f"slide_{slide_number}_shape_{shape_index}",
                            "conversion_method": "image_fallback"
                        })

            except Exception as e:
                print(f"处理幻灯片 {slide_number} 形状 {shape_index} 时出错: {e}")
                continue

        return equations

    def _get_shape_xml(self, shape):
        """获取形状的XML内容"""
        try:
            # 尝试获取形状的内部XML
            if hasattr(shape, '_element'):
                return ET.tostring(shape._element, encoding='unicode')
        except Exception as e:
            print(f"获取形状XML失败: {e}")
        return None

    def _extract_omml_from_xml(self, xml_string):
        """从XML中提取OMML内容"""
        try:
            # 查找OMML数学标签
            omml_patterns = [
                r'<m:oMath[^>]*>.*?</m:oMath>',
                r'<m:oMathPara[^>]*>.*?</m:oMathPara>',
                r'<math[^>]*>.*?</math>'  # 也检查标准MathML
            ]

            for pattern in omml_patterns:
                matches = re.findall(pattern, xml_string, re.DOTALL | re.IGNORECASE)
                if matches:
                    return matches[0]

        except Exception as e:
            print(f"提取OMML失败: {e}")
        return None

    def _convert_omml_to_latex(self, omml_content):
        """将OMML转换为LaTeX"""
        try:
            # 方法1: 尝试使用pandoc
            latex_result = self._convert_via_pandoc(omml_content)
            if latex_result:
                return latex_result

            # 方法2: 简单的文本替换作为备选方案
            latex_result = self._simple_omml_to_latex(omml_content)
            if latex_result:
                return latex_result

        except Exception as e:
            print(f"OMML转LaTeX失败: {e}")

        return None

    def _convert_via_pandoc(self, omml_content):
        """使用pandoc转换OMML到LaTeX"""
        try:
            # 检查pandoc是否可用
            subprocess.run(['pandoc', '--version'],
                           capture_output=True, check=True)

            # 创建临时文件
            with tempfile.NamedTemporaryFile(mode='w', suffix='.xml', delete=False) as temp_file:
                temp_file.write(f'<root>{omml_content}</root>')
                temp_file_path = temp_file.name

            try:
                # 使用pandoc转换
                result = subprocess.run([
                    'pandoc',
                    '-f', 'docx',
                    '-t', 'latex',
                    temp_file_path
                ], capture_output=True, text=True, check=True)

                return result.stdout.strip()

            finally:
                os.unlink(temp_file_path)

        except (subprocess.CalledProcessError, FileNotFoundError):
            print("Pandoc不可用，跳过pandoc转换")
        except Exception as e:
            print(f"Pandoc转换失败: {e}")

        return None

    def _simple_omml_to_latex(self, omml_content):
        """简单的OMML到LaTeX转换（基本文本替换）"""
        try:
            # 移除XML标签，提取纯文本
            text_content = re.sub(r'<[^>]+>', '', omml_content)
            text_content = text_content.strip()

            if not text_content:
                return None

            # 基本的数学符号替换
            replacements = {
                '≈': r'\approx',
                '≠': r'\neq',
                '≤': r'\leq',
                '≥': r'\geq',
                '∞': r'\infty',
                'α': r'\alpha',
                'β': r'\beta',
                'γ': r'\gamma',
                'δ': r'\delta',
                'θ': r'\theta',
                'λ': r'\lambda',
                'μ': r'\mu',
                'π': r'\pi',
                'σ': r'\sigma',
                'φ': r'\phi',
                'ω': r'\omega',
                '∑': r'\sum',
                '∫': r'\int',
                '√': r'\sqrt',
                '±': r'\pm',
                '×': r'\times',
                '÷': r'\div'
            }

            for symbol, latex in replacements.items():
                text_content = text_content.replace(symbol, latex)

            # 包装在数学环境中
            return f"${text_content}$"

        except Exception as e:
            print(f"简单转换失败: {e}")

        return None

    def _is_potential_formula_image(self, shape):
        """判断形状是否可能是公式图片"""
        try:
            # 检查是否为图片类型
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                return True

            # 检查是否为包含复杂路径的形状（可能是矢量公式）
            if shape.shape_type in [MSO_SHAPE_TYPE.FREEFORM, MSO_SHAPE_TYPE.AUTO_SHAPE]:
                return True

            # 检查形状大小（小的形状可能是公式）
            if hasattr(shape, 'width') and hasattr(shape, 'height'):
                # 假设公式通常比较小（宽度和高度都小于某个阈值）
                max_formula_size = 200000  # EMU单位
                if shape.width < max_formula_size and shape.height < max_formula_size:
                    return True

        except Exception as e:
            print(f"检查潜在公式图片失败: {e}")

        return False

    def _process_formula_image(self, shape, slide_number, shape_index):
        """处理公式图片"""
        try:
            # 如果是图片类型，尝试导出图片
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_filename = f"formula_slide_{slide_number}_shape_{shape_index}.png"
                image_path = os.path.join(self.output_img_dir, image_filename)

                # 这里需要实现图片导出逻辑
                # 由于python-pptx的限制，可能需要使用其他方法
                print(f"识别到公式图片，但需要额外的导出逻辑: {image_filename}")

                return image_path

        except Exception as e:
            print(f"处理公式图片失败: {e}")

        return None

    def process_pptx_file_advanced(self, file_path):
        """
        高级PPTX处理：结合传统方法和ZIP解析

        Args:
            file_path: PPTX文件路径
        """
        print(f"开始高级PPTX处理: {file_path}")
        base_filename = os.path.splitext(os.path.basename(file_path))[0]

        # 重置结果记录，为当前PPTX文件单独记录
        current_results = []
        original_results = self.preprocessor.results
        self.preprocessor.results = current_results

        try:
            # 1. 使用传统python-pptx方法处理文本和表格
            self._process_text_and_tables_traditional(file_path, base_filename)

            # 2. 使用ZIP方法提取所有图片
            slide_image_mapping = self.extract_all_images_via_zip(file_path)

            # 3. 生成PPTX专用元数据
            self._save_pptx_metadata(file_path, base_filename, slide_image_mapping)

            print(f"高级PPTX处理完成: {file_path}")
            print(f"PPTX元数据已保存: output/{base_filename}_pptx_metadata.json")

        except Exception as e:
            print(f"高级PPTX处理失败: {e}")
            import traceback
            traceback.print_exc()
        finally:
            # 恢复原始结果列表并合并当前结果
            self.preprocessor.results = original_results
            self.preprocessor.results.extend(current_results)

    def _process_text_and_tables_traditional(self, file_path, base_filename):
        """使用传统python-pptx方法处理文本和表格"""
        prs = Presentation(file_path)

        for slide_index, slide in enumerate(prs.slides, start=1):
            # 处理公式
            equations = self.extract_and_convert_equations(slide, slide_index)
            if equations:
                print(f"在幻灯片 {slide_index} 中找到 {len(equations)} 个公式")

                # 保存公式信息到JSON文件
                formulas_json_path = f"output/formulas/{base_filename}_slide_{slide_index}_formulas.json"
                os.makedirs("output/formulas", exist_ok=True)

                formulas_output = {
                    "slide_number": slide_index,
                    "source_file": base_filename,
                    "equations_count": len(equations),
                    "equations": equations,
                    "processing_date": str(pd.Timestamp.now())
                }

                with open(formulas_json_path, "w", encoding="utf-8") as f:
                    json.dump(formulas_output, f, ensure_ascii=False, indent=2)

                print(f"公式信息已保存到: {formulas_json_path}")

            # 提取文本
            slide_text_items = []
            for shape in slide.shapes:
                if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                    text_content = []
                    for paragraph in shape.text_frame.paragraphs:
                        runs_text = ''.join(run.text for run in paragraph.runs)
                        text_content.append(runs_text if runs_text else paragraph.text)
                    final_text = "\n".join([t for t in text_content if t is not None])
                    if final_text.strip():
                        slide_text_items.append(final_text.strip())

            if slide_text_items:
                text_output = {
                    "type": "ppt_text",
                    "page": slide_index,
                    "source": f"{base_filename}_slide_{slide_index}",
                    "raw_text": "\n\n".join(slide_text_items)
                }
                text_json_path = f"{self.output_text_dir}/{base_filename}_slide_{slide_index}.json"
                with open(text_json_path, "w", encoding="utf-8") as f:
                    json.dump(text_output, f, ensure_ascii=False, indent=2)
                self.preprocessor.results.append({
                    "type": "ppt_text",
                    "page": slide_index,
                    "file": text_json_path
                })

            # 提取表格（优化版本）
            table_counter = 0
            for shape in slide.shapes:
                if hasattr(shape, "has_table") and shape.has_table:
                    table_counter += 1
                    table = shape.table

                    # 获取表格位置信息（增强功能）
                    table_position = {
                        "left": float(shape.left.inches) if shape.left else 0,
                        "top": float(shape.top.inches) if shape.top else 0,
                        "width": float(shape.width.inches) if shape.width else 0,
                        "height": float(shape.height.inches) if shape.height else 0
                    }

                    # 提取表格数据
                    data_matrix = []
                    for row in table.rows:
                        row_values = []
                        for cell in row.cells:
                            # 优化：使用 text_frame.text 获取纯文本
                            try:
                                if cell.text_frame and cell.text_frame.text:
                                    cell_text = cell.text_frame.text.strip()
                                else:
                                    cell_text = cell.text.strip() if cell.text else ""
                            except Exception as e:
                                print(f"提取单元格文本失败: {e}")
                                cell_text = ""
                            row_values.append(cell_text)
                        data_matrix.append(row_values)

                    # 检查是否有有效数据
                    if data_matrix and any(any(cell for cell in row) for row in data_matrix):
                        # 使用pandas DataFrame保存为CSV
                        df = pd.DataFrame(data_matrix)
                        csv_path = f"{self.output_table_dir}/{base_filename}_slide_{slide_index}_table_{table_counter}.csv"
                        df.to_csv(csv_path, index=False, header=False, encoding="utf-8")

                        # 创建表格JSON元数据文件
                        table_metadata = {
                            "type": "ppt_table",
                            "source": f"{base_filename}_slide_{slide_index}",
                            "slide_number": slide_index,
                            "table_index": table_counter,
                            "dimensions": {
                                "rows": len(data_matrix),
                                "columns": len(data_matrix[0]) if data_matrix else 0
                            },
                            "position": table_position,
                            "data_preview": data_matrix[:3] if len(data_matrix) > 0 else [],  # 前3行预览
                            "csv_file": csv_path
                        }

                        # 保存表格元数据JSON
                        json_path = f"{self.output_table_dir}/{base_filename}_slide_{slide_index}_table_{table_counter}.json"
                        with open(json_path, "w", encoding="utf-8") as f:
                            json.dump(table_metadata, f, ensure_ascii=False, indent=2)

                        # 记录到主结果中（增强元数据）
                        self.preprocessor.results.append({
                            "type": "ppt_table",
                            "page": slide_index,
                            "table_index": table_counter,
                            "file": csv_path,
                            "metadata_file": json_path,
                            "dimensions": {
                                "rows": len(data_matrix),
                                "columns": len(data_matrix[0]) if data_matrix else 0
                            },
                            "position": table_position,
                            "extraction_method": "python_pptx_optimized"
                        })

                        print(
                            f"✓ 提取表格 {table_counter}: {len(data_matrix)}行 x {len(data_matrix[0]) if data_matrix else 0}列")
                        print(f"  位置: left={table_position['left']:.2f}in, top={table_position['top']:.2f}in")
                    else:
                        print(f"⚠ 跳过空表格 {table_counter}")

    def _save_pptx_metadata(self, file_path, base_filename, slide_image_mapping):
        """
        保存PPTX专用元数据文件

        Args:
            file_path: 原始PPTX文件路径
            base_filename: 基础文件名
            slide_image_mapping: 幻灯片到图片的映射关系
        """
        import datetime

        # 统计各类型文件
        text_files = [r for r in self.preprocessor.results if r["type"] == "ppt_text"]
        table_files = [r for r in self.preprocessor.results if r["type"] == "ppt_table"]
        image_files_traditional = [r for r in self.preprocessor.results if r["type"] == "ppt_image"]
        image_files_zip = [r for r in self.preprocessor.results if r["type"] == "ppt_image_zip"]

        # 计算幻灯片统计
        total_slides = len(set(r["page"] for r in self.preprocessor.results if "page" in r))
        slides_with_images = len(slide_image_mapping)
        total_images_zip = sum(len(images) for images in slide_image_mapping.values())

        # 计算表格统计
        total_tables = len(table_files)
        total_table_rows = sum(r.get("dimensions", {}).get("rows", 0) for r in table_files)
        total_table_columns = sum(r.get("dimensions", {}).get("columns", 0) for r in table_files)
        table_positions = [r.get("position", {}) for r in table_files if r.get("position")]

        # 构建元数据
        metadata = {
            "processing_date": datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
            "source_file": base_filename,
            "source_path": file_path,
            "file_type": "PPTX",
            "processing_method": "advanced_zip_xml_parsing",
            "statistics": {
                "total_slides": total_slides,
                "slides_with_text": len(text_files),
                "slides_with_tables": len([r for r in table_files if r.get("dimensions", {}).get("rows", 0) > 0]),
                "total_tables": total_tables,
                "total_table_rows": total_table_rows,
                "total_table_columns": total_table_columns,
                "slides_with_images": slides_with_images,
                "total_images_extracted": len(image_files_traditional) + len(image_files_zip),
                "images_via_traditional": len(image_files_traditional),
                "images_via_zip_parsing": len(image_files_zip),
                "total_images_in_media": total_images_zip
            },
            "slide_image_mapping": slide_image_mapping,
            "files": {
                "text_files": text_files,
                "table_files": table_files,
                "image_files_traditional": image_files_traditional,
                "image_files_zip": image_files_zip
            },
            "table_analysis": {
                "positions": table_positions,
                "position_stats": {
                    "avg_left": sum(p.get("left", 0) for p in table_positions) / len(
                        table_positions) if table_positions else 0,
                    "avg_top": sum(p.get("top", 0) for p in table_positions) / len(
                        table_positions) if table_positions else 0,
                    "avg_width": sum(p.get("width", 0) for p in table_positions) / len(
                        table_positions) if table_positions else 0,
                    "avg_height": sum(p.get("height", 0) for p in table_positions) / len(
                        table_positions) if table_positions else 0
                }
            },
            "processing_info": {
                "extraction_methods": ["python-pptx", "zip_xml_parsing"],
                "image_formats_supported": ["PNG", "WMF", "EMF", "JPEG"],
                "table_extraction_enhanced": True,
                "table_position_tracking": True,
                "clip_descriptions_generated": True,
                "output_format": "JSON/CSV"
            }
        }

        # 保存元数据
        metadata_path = f"output/{base_filename}_pptx_metadata.json"
        with open(metadata_path, "w", encoding="utf-8") as f:
            json.dump(metadata, f, ensure_ascii=False, indent=2)

        return metadata_path


def process_pptx_file_advanced(preprocessor, file_path, fast_mode=False):
    """
    高级PPTX处理的入口函数

    Args:
        preprocessor: MultimodalPreprocessor实例
        file_path: PPTX文件路径
        fast_mode: 快速模式，跳过耗时处理
    """
    processor = AdvancedPPTProcessor(preprocessor, fast_mode=fast_mode)
    processor.process_pptx_file_advanced(file_path)


import sys
import os

# 获取当前文件的绝对路径
current_dir = os.path.dirname(os.path.abspath(__file__))
parent_dir = os.path.dirname(current_dir)

# 添加路径
sys.path.append(parent_dir)  # 项目根目录
sys.path.append(current_dir)  # src目录
import json
import fitz  # PyMuPDF
import torch
import numpy as np
from PIL import Image, ImageEnhance
from transformers import CLIPProcessor, CLIPModel
import datetime
import pandas as pd
import cv2
import easyocr
import re
import pdfplumber
import camelot
import csv

'''
try:
    # 懒加载高级PPTX处理器（若不可用则忽略）
    from advanced_pptx_processor import process_pptx_file_advanced
except Exception:
    process_pptx_file_advanced = None
    '''
import os
import json
import zipfile
import tempfile
import shutil
import xml.etree.ElementTree as ET

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import pandas as pd
import re
import subprocess


class AdvancedPPTProcessor:
    def __init__(self, preprocessor, fast_mode=False):
        """
        初始化高级PPTX处理器

        Args:
            preprocessor: MultimodalPreprocessor实例，用于重用输出目录与结果记录
            fast_mode: 快速模式，跳过耗时的CLIP描述生成
        """
        self.preprocessor = preprocessor
        self.fast_mode = fast_mode
        self.output_text_dir = "output/text"
        self.output_table_dir = "output/tables"
        self.output_img_dir = "output/images"

        # 确保输出目录存在
        os.makedirs(self.output_text_dir, exist_ok=True)
        os.makedirs(self.output_table_dir, exist_ok=True)
        os.makedirs(self.output_img_dir, exist_ok=True)

    def extract_all_images_via_zip(self, file_path):
        """
        通过ZIP解压和XML解析提取PPTX中的所有图片

        Args:
            file_path: PPTX文件路径

        Returns:
            dict: 包含幻灯片到图片映射关系的字典
        """
        print(f"开始通过ZIP方式提取图片: {file_path}")

        base_filename = os.path.splitext(os.path.basename(file_path))[0]
        slide_image_mapping = {}

        # 创建临时目录
        with tempfile.TemporaryDirectory() as temp_dir:
            try:
                # 1. 解压PPTX文件
                print("正在解压PPTX文件...")
                with zipfile.ZipFile(file_path, 'r') as zip_ref:
                    zip_ref.extractall(temp_dir)

                # 2. 找到媒体目录
                media_dir = os.path.join(temp_dir, "ppt", "media")
                slides_dir = os.path.join(temp_dir, "ppt", "slides")
                rels_dir = os.path.join(temp_dir, "ppt", "slides", "_rels")

                if not os.path.exists(media_dir):
                    print("未找到media目录，可能没有图片")
                    return slide_image_mapping

                print(f"找到media目录: {media_dir}")
                print(f"媒体文件: {os.listdir(media_dir)}")

                # 3. 遍历所有幻灯片XML文件
                if os.path.exists(slides_dir):
                    for slide_file in os.listdir(slides_dir):
                        if slide_file.startswith("slide") and slide_file.endswith(".xml"):
                            slide_num = self._extract_slide_number(slide_file)
                            if slide_num is None:
                                continue

                            print(f"处理幻灯片 {slide_num}: {slide_file}")

                            # 解析幻灯片XML获取图片关系ID
                            slide_xml_path = os.path.join(slides_dir, slide_file)
                            image_rids = self._parse_slide_xml_for_images(slide_xml_path)

                            if image_rids:
                                print(f"幻灯片 {slide_num} 中找到图片关系ID: {image_rids}")

                                # 解析关系文件获取实际文件名
                                rels_file = slide_file + ".rels"
                                rels_path = os.path.join(rels_dir, rels_file)

                                if os.path.exists(rels_path):
                                    image_files = self._parse_rels_file(rels_path, image_rids)

                                    if image_files:
                                        slide_image_mapping[slide_num] = image_files
                                        print(f"幻灯片 {slide_num} 映射到图片: {image_files}")

                                        # 复制图片到输出目录
                                        self._copy_images_to_output(media_dir, image_files,
                                                                    base_filename, slide_num)

                print(f"图片提取完成，映射关系: {slide_image_mapping}")

            except Exception as e:
                print(f"ZIP方式图片提取失败: {e}")
                import traceback
                traceback.print_exc()

        return slide_image_mapping

    def _extract_slide_number(self, slide_filename):
        """从幻灯片文件名中提取编号"""
        try:
            # slide1.xml -> 1
            import re
            match = re.search(r'slide(\d+)\.xml', slide_filename)
            if match:
                return int(match.group(1))
        except Exception:
            pass
        return None

    def _parse_slide_xml_for_images(self, slide_xml_path):
        """
        解析幻灯片XML文件，查找图片引用

        Args:
            slide_xml_path: 幻灯片XML文件路径

        Returns:
            list: 图片关系ID列表
        """
        image_rids = []

        try:
            tree = ET.parse(slide_xml_path)
            root = tree.getroot()

            # 定义命名空间
            namespaces = {
                'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
                'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
                'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'
            }

            # 查找所有a:blip元素（图片引用）
            blip_elements = root.findall('.//a:blip', namespaces)

            for blip in blip_elements:
                embed_attr = blip.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
                if embed_attr:
                    image_rids.append(embed_attr)
                    print(f"找到图片引用ID: {embed_attr}")

        except Exception as e:
            print(f"解析幻灯片XML失败 {slide_xml_path}: {e}")

        return image_rids

    def _parse_rels_file(self, rels_path, image_rids):
        """
        解析关系文件，获取关系ID到文件名的映射

        Args:
            rels_path: 关系文件路径
            image_rids: 图片关系ID列表

        Returns:
            list: 对应的图片文件名列表
        """
        image_files = []

        try:
            tree = ET.parse(rels_path)
            root = tree.getroot()

            # 定义命名空间
            namespaces = {
                'rel': 'http://schemas.openxmlformats.org/package/2006/relationships'
            }

            # 查找所有关系
            for relationship in root.findall('.//rel:Relationship', namespaces):
                rel_id = relationship.get('Id')
                target = relationship.get('Target')
                rel_type = relationship.get('Type')

                # 检查是否是图片关系
                if (rel_id in image_rids and
                        target and
                        rel_type and
                        'image' in rel_type.lower()):
                    # 提取文件名 (../media/image1.png -> image1.png)
                    filename = os.path.basename(target)
                    image_files.append(filename)
                    print(f"关系映射: {rel_id} -> {filename}")

        except Exception as e:
            print(f"解析关系文件失败 {rels_path}: {e}")

        return image_files

    def _copy_images_to_output(self, media_dir, image_files, base_filename, slide_num):
        """
        将图片复制到输出目录并生成描述

        Args:
            media_dir: 媒体文件源目录
            image_files: 图片文件名列表
            base_filename: 基础文件名
            slide_num: 幻灯片编号
        """
        for idx, image_file in enumerate(image_files, 1):
            try:
                source_path = os.path.join(media_dir, image_file)

                if os.path.exists(source_path):
                    # 生成输出文件名
                    file_ext = os.path.splitext(image_file)[1]
                    output_filename = f"{base_filename}_slide_{slide_num}_img_{idx}_zip{file_ext}"
                    output_path = os.path.join(self.output_img_dir, output_filename)

                    # 复制图片
                    shutil.copy2(source_path, output_path)
                    print(f"复制图片: {source_path} -> {output_path}")

                    # 生成CLIP描述（根据模式决定是否生成）
                    desc_path = None
                    if not self.fast_mode:
                        try:
                            desc_path = self.preprocessor.clip_generate_description(output_path)
                        except Exception as e:
                            print(f"生成CLIP描述失败，跳过: {e}")
                    else:
                        print("快速模式：跳过CLIP描述生成")

                    # 记录到结果中
                    self.preprocessor.results.append({
                        "type": "ppt_image_zip",
                        "page": slide_num,
                        "file": output_path,
                        "description_file": desc_path,
                        "extraction_method": "zip_xml_parsing",
                        "original_filename": image_file
                    })

                else:
                    print(f"源图片文件不存在: {source_path}")

            except Exception as e:
                print(f"复制图片失败 {image_file}: {e}")

    def extract_and_convert_equations(self, slide, slide_number):
        """
        处理幻灯片中的公式

        Args:
            slide: python-pptx的Slide对象
            slide_number: 幻灯片编号

        Returns:
            list: 包含公式信息的列表
        """
        equations = []

        for shape_index, shape in enumerate(slide.shapes):
            try:
                # 检查形状是否包含文本框
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    # 获取形状的XML内容
                    shape_xml = self._get_shape_xml(shape)
                    if shape_xml:
                        # 检查是否包含OMML公式标签
                        omml_content = self._extract_omml_from_xml(shape_xml)
                        if omml_content:
                            print(f"在幻灯片 {slide_number} 形状 {shape_index} 中发现OMML公式")

                            # 尝试转换OMML到LaTeX
                            latex_content = self._convert_omml_to_latex(omml_content)

                            equation_info = {
                                "slide_number": slide_number,
                                "shape_index": shape_index,
                                "type": "omml_formula",
                                "original_omml": omml_content[:500] + "..." if len(
                                    omml_content) > 500 else omml_content,
                                "latex": latex_content,
                                "conversion_success": latex_content is not None
                            }

                            equations.append(equation_info)

                            # 添加到结果中
                            self.preprocessor.results.append({
                                "type": "formula",
                                "page": slide_number,
                                "formula_type": "omml",
                                "latex": latex_content,
                                "source": f"slide_{slide_number}_shape_{shape_index}",
                                "conversion_method": "omml_to_latex"
                            })

                            continue

                # 如果没有找到OMML，检查是否为可能的公式图片
                if self._is_potential_formula_image(shape):
                    print(f"在幻灯片 {slide_number} 形状 {shape_index} 中发现潜在公式图片")

                    # 使用图片处理流程处理公式图片
                    formula_image_path = self._process_formula_image(shape, slide_number, shape_index)

                    if formula_image_path:
                        equation_info = {
                            "slide_number": slide_number,
                            "shape_index": shape_index,
                            "type": "formula_image",
                            "image_path": formula_image_path,
                            "latex": None,
                            "conversion_success": False
                        }

                        equations.append(equation_info)

                        # 添加到结果中
                        self.preprocessor.results.append({
                            "type": "formula",
                            "page": slide_number,
                            "formula_type": "image",
                            "image_path": formula_image_path,
                            "source": f"slide_{slide_number}_shape_{shape_index}",
                            "conversion_method": "image_fallback"
                        })

            except Exception as e:
                print(f"处理幻灯片 {slide_number} 形状 {shape_index} 时出错: {e}")
                continue

        return equations

    def _get_shape_xml(self, shape):
        """获取形状的XML内容"""
        try:
            # 尝试获取形状的内部XML
            if hasattr(shape, '_element'):
                return ET.tostring(shape._element, encoding='unicode')
        except Exception as e:
            print(f"获取形状XML失败: {e}")
        return None

    def _extract_omml_from_xml(self, xml_string):
        """从XML中提取OMML内容"""
        try:
            # 查找OMML数学标签
            omml_patterns = [
                r'<m:oMath[^>]*>.*?</m:oMath>',
                r'<m:oMathPara[^>]*>.*?</m:oMathPara>',
                r'<math[^>]*>.*?</math>'  # 也检查标准MathML
            ]

            for pattern in omml_patterns:
                matches = re.findall(pattern, xml_string, re.DOTALL | re.IGNORECASE)
                if matches:
                    return matches[0]

        except Exception as e:
            print(f"提取OMML失败: {e}")
        return None

    def _convert_omml_to_latex(self, omml_content):
        """将OMML转换为LaTeX"""
        try:
            # 方法1: 尝试使用pandoc
            latex_result = self._convert_via_pandoc(omml_content)
            if latex_result:
                return latex_result

            # 方法2: 简单的文本替换作为备选方案
            latex_result = self._simple_omml_to_latex(omml_content)
            if latex_result:
                return latex_result

        except Exception as e:
            print(f"OMML转LaTeX失败: {e}")

        return None

    def _convert_via_pandoc(self, omml_content):
        """使用pandoc转换OMML到LaTeX"""
        try:
            # 检查pandoc是否可用
            subprocess.run(['pandoc', '--version'],
                           capture_output=True, check=True)

            # 创建临时文件
            with tempfile.NamedTemporaryFile(mode='w', suffix='.xml', delete=False) as temp_file:
                temp_file.write(f'<root>{omml_content}</root>')
                temp_file_path = temp_file.name

            try:
                # 使用pandoc转换
                result = subprocess.run([
                    'pandoc',
                    '-f', 'docx',
                    '-t', 'latex',
                    temp_file_path
                ], capture_output=True, text=True, check=True)

                return result.stdout.strip()

            finally:
                os.unlink(temp_file_path)

        except (subprocess.CalledProcessError, FileNotFoundError):
            print("Pandoc不可用，跳过pandoc转换")
        except Exception as e:
            print(f"Pandoc转换失败: {e}")

        return None

    def _simple_omml_to_latex(self, omml_content):
        """简单的OMML到LaTeX转换（基本文本替换）"""
        try:
            # 移除XML标签，提取纯文本
            text_content = re.sub(r'<[^>]+>', '', omml_content)
            text_content = text_content.strip()

            if not text_content:
                return None

            # 基本的数学符号替换
            replacements = {
                '≈': r'\approx',
                '≠': r'\neq',
                '≤': r'\leq',
                '≥': r'\geq',
                '∞': r'\infty',
                'α': r'\alpha',
                'β': r'\beta',
                'γ': r'\gamma',
                'δ': r'\delta',
                'θ': r'\theta',
                'λ': r'\lambda',
                'μ': r'\mu',
                'π': r'\pi',
                'σ': r'\sigma',
                'φ': r'\phi',
                'ω': r'\omega',
                '∑': r'\sum',
                '∫': r'\int',
                '√': r'\sqrt',
                '±': r'\pm',
                '×': r'\times',
                '÷': r'\div'
            }

            for symbol, latex in replacements.items():
                text_content = text_content.replace(symbol, latex)

            # 包装在数学环境中
            return f"${text_content}$"

        except Exception as e:
            print(f"简单转换失败: {e}")

        return None

    def _is_potential_formula_image(self, shape):
        """判断形状是否可能是公式图片"""
        try:
            # 检查是否为图片类型
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                return True

            # 检查是否为包含复杂路径的形状（可能是矢量公式）
            if shape.shape_type in [MSO_SHAPE_TYPE.FREEFORM, MSO_SHAPE_TYPE.AUTO_SHAPE]:
                return True

            # 检查形状大小（小的形状可能是公式）
            if hasattr(shape, 'width') and hasattr(shape, 'height'):
                # 假设公式通常比较小（宽度和高度都小于某个阈值）
                max_formula_size = 200000  # EMU单位
                if shape.width < max_formula_size and shape.height < max_formula_size:
                    return True

        except Exception as e:
            print(f"检查潜在公式图片失败: {e}")

        return False

    def _process_formula_image(self, shape, slide_number, shape_index):
        """处理公式图片"""
        try:
            # 如果是图片类型，尝试导出图片
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_filename = f"formula_slide_{slide_number}_shape_{shape_index}.png"
                image_path = os.path.join(self.output_img_dir, image_filename)

                # 这里需要实现图片导出逻辑
                # 由于python-pptx的限制，可能需要使用其他方法
                print(f"识别到公式图片，但需要额外的导出逻辑: {image_filename}")

                return image_path

        except Exception as e:
            print(f"处理公式图片失败: {e}")

        return None

    def process_pptx_file_advanced(self, file_path):
        """
        高级PPTX处理：结合传统方法和ZIP解析

        Args:
            file_path: PPTX文件路径
        """
        print(f"开始高级PPTX处理: {file_path}")
        base_filename = os.path.splitext(os.path.basename(file_path))[0]

        # 重置结果记录，为当前PPTX文件单独记录
        current_results = []
        original_results = self.preprocessor.results
        self.preprocessor.results = current_results

        try:
            # 1. 使用传统python-pptx方法处理文本和表格
            self._process_text_and_tables_traditional(file_path, base_filename)

            # 2. 使用ZIP方法提取所有图片
            slide_image_mapping = self.extract_all_images_via_zip(file_path)

            # 3. 生成PPTX专用元数据
            self._save_pptx_metadata(file_path, base_filename, slide_image_mapping)

            print(f"高级PPTX处理完成: {file_path}")
            print(f"PPTX元数据已保存: output/{base_filename}_pptx_metadata.json")

        except Exception as e:
            print(f"高级PPTX处理失败: {e}")
            import traceback
            traceback.print_exc()
        finally:
            # 恢复原始结果列表并合并当前结果
            self.preprocessor.results = original_results
            self.preprocessor.results.extend(current_results)

    def _process_text_and_tables_traditional(self, file_path, base_filename):
        """使用传统python-pptx方法处理文本和表格"""
        prs = Presentation(file_path)

        for slide_index, slide in enumerate(prs.slides, start=1):
            # 处理公式
            equations = self.extract_and_convert_equations(slide, slide_index)
            if equations:
                print(f"在幻灯片 {slide_index} 中找到 {len(equations)} 个公式")

                # 保存公式信息到JSON文件
                formulas_json_path = f"output/formulas/{base_filename}_slide_{slide_index}_formulas.json"
                os.makedirs("output/formulas", exist_ok=True)

                formulas_output = {
                    "slide_number": slide_index,
                    "source_file": base_filename,
                    "equations_count": len(equations),
                    "equations": equations,
                    "processing_date": str(pd.Timestamp.now())
                }

                with open(formulas_json_path, "w", encoding="utf-8") as f:
                    json.dump(formulas_output, f, ensure_ascii=False, indent=2)

                print(f"公式信息已保存到: {formulas_json_path}")

            # 提取文本
            slide_text_items = []
            for shape in slide.shapes:
                if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                    text_content = []
                    for paragraph in shape.text_frame.paragraphs:
                        runs_text = ''.join(run.text for run in paragraph.runs)
                        text_content.append(runs_text if runs_text else paragraph.text)
                    final_text = "\n".join([t for t in text_content if t is not None])
                    if final_text.strip():
                        slide_text_items.append(final_text.strip())

            if slide_text_items:
                text_output = {
                    "type": "ppt_text",
                    "page": slide_index,
                    "source": f"{base_filename}_slide_{slide_index}",
                    "raw_text": "\n\n".join(slide_text_items)
                }
                text_json_path = f"{self.output_text_dir}/{base_filename}_slide_{slide_index}.json"
                with open(text_json_path, "w", encoding="utf-8") as f:
                    json.dump(text_output, f, ensure_ascii=False, indent=2)
                self.preprocessor.results.append({
                    "type": "ppt_text",
                    "page": slide_index,
                    "file": text_json_path
                })

            # 提取表格（优化版本）
            table_counter = 0
            for shape in slide.shapes:
                if hasattr(shape, "has_table") and shape.has_table:
                    table_counter += 1
                    table = shape.table

                    # 获取表格位置信息（增强功能）
                    table_position = {
                        "left": float(shape.left.inches) if shape.left else 0,
                        "top": float(shape.top.inches) if shape.top else 0,
                        "width": float(shape.width.inches) if shape.width else 0,
                        "height": float(shape.height.inches) if shape.height else 0
                    }

                    # 提取表格数据
                    data_matrix = []
                    for row in table.rows:
                        row_values = []
                        for cell in row.cells:
                            # 优化：使用 text_frame.text 获取纯文本
                            try:
                                if cell.text_frame and cell.text_frame.text:
                                    cell_text = cell.text_frame.text.strip()
                                else:
                                    cell_text = cell.text.strip() if cell.text else ""
                            except Exception as e:
                                print(f"提取单元格文本失败: {e}")
                                cell_text = ""
                            row_values.append(cell_text)
                        data_matrix.append(row_values)

                    # 检查是否有有效数据
                    if data_matrix and any(any(cell for cell in row) for row in data_matrix):
                        # 使用pandas DataFrame保存为CSV
                        df = pd.DataFrame(data_matrix)
                        csv_path = f"{self.output_table_dir}/{base_filename}_slide_{slide_index}_table_{table_counter}.csv"
                        df.to_csv(csv_path, index=False, header=False, encoding="utf-8")

                        # 创建表格JSON元数据文件
                        table_metadata = {
                            "type": "ppt_table",
                            "source": f"{base_filename}_slide_{slide_index}",
                            "slide_number": slide_index,
                            "table_index": table_counter,
                            "dimensions": {
                                "rows": len(data_matrix),
                                "columns": len(data_matrix[0]) if data_matrix else 0
                            },
                            "position": table_position,
                            "data_preview": data_matrix[:3] if len(data_matrix) > 0 else [],  # 前3行预览
                            "csv_file": csv_path
                        }

                        # 保存表格元数据JSON
                        json_path = f"{self.output_table_dir}/{base_filename}_slide_{slide_index}_table_{table_counter}.json"
                        with open(json_path, "w", encoding="utf-8") as f:
                            json.dump(table_metadata, f, ensure_ascii=False, indent=2)

                        # 记录到主结果中（增强元数据）
                        self.preprocessor.results.append({
                            "type": "ppt_table",
                            "page": slide_index,
                            "table_index": table_counter,
                            "file": csv_path,
                            "metadata_file": json_path,
                            "dimensions": {
                                "rows": len(data_matrix),
                                "columns": len(data_matrix[0]) if data_matrix else 0
                            },
                            "position": table_position,
                            "extraction_method": "python_pptx_optimized"
                        })

                        print(
                            f"✓ 提取表格 {table_counter}: {len(data_matrix)}行 x {len(data_matrix[0]) if data_matrix else 0}列")
                        print(f"  位置: left={table_position['left']:.2f}in, top={table_position['top']:.2f}in")
                    else:
                        print(f"⚠ 跳过空表格 {table_counter}")

    def _save_pptx_metadata(self, file_path, base_filename, slide_image_mapping):
        """
        保存PPTX专用元数据文件

        Args:
            file_path: 原始PPTX文件路径
            base_filename: 基础文件名
            slide_image_mapping: 幻灯片到图片的映射关系
        """
        import datetime

        # 统计各类型文件
        text_files = [r for r in self.preprocessor.results if r["type"] == "ppt_text"]
        table_files = [r for r in self.preprocessor.results if r["type"] == "ppt_table"]
        image_files_traditional = [r for r in self.preprocessor.results if r["type"] == "ppt_image"]
        image_files_zip = [r for r in self.preprocessor.results if r["type"] == "ppt_image_zip"]

        # 计算幻灯片统计
        total_slides = len(set(r["page"] for r in self.preprocessor.results if "page" in r))
        slides_with_images = len(slide_image_mapping)
        total_images_zip = sum(len(images) for images in slide_image_mapping.values())

        # 计算表格统计
        total_tables = len(table_files)
        total_table_rows = sum(r.get("dimensions", {}).get("rows", 0) for r in table_files)
        total_table_columns = sum(r.get("dimensions", {}).get("columns", 0) for r in table_files)
        table_positions = [r.get("position", {}) for r in table_files if r.get("position")]

        # 构建元数据
        metadata = {
            "processing_date": datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
            "source_file": base_filename,
            "source_path": file_path,
            "file_type": "PPTX",
            "processing_method": "advanced_zip_xml_parsing",
            "statistics": {
                "total_slides": total_slides,
                "slides_with_text": len(text_files),
                "slides_with_tables": len([r for r in table_files if r.get("dimensions", {}).get("rows", 0) > 0]),
                "total_tables": total_tables,
                "total_table_rows": total_table_rows,
                "total_table_columns": total_table_columns,
                "slides_with_images": slides_with_images,
                "total_images_extracted": len(image_files_traditional) + len(image_files_zip),
                "images_via_traditional": len(image_files_traditional),
                "images_via_zip_parsing": len(image_files_zip),
                "total_images_in_media": total_images_zip
            },
            "slide_image_mapping": slide_image_mapping,
            "files": {
                "text_files": text_files,
                "table_files": table_files,
                "image_files_traditional": image_files_traditional,
                "image_files_zip": image_files_zip
            },
            "table_analysis": {
                "positions": table_positions,
                "position_stats": {
                    "avg_left": sum(p.get("left", 0) for p in table_positions) / len(
                        table_positions) if table_positions else 0,
                    "avg_top": sum(p.get("top", 0) for p in table_positions) / len(
                        table_positions) if table_positions else 0,
                    "avg_width": sum(p.get("width", 0) for p in table_positions) / len(
                        table_positions) if table_positions else 0,
                    "avg_height": sum(p.get("height", 0) for p in table_positions) / len(
                        table_positions) if table_positions else 0
                }
            },
            "processing_info": {
                "extraction_methods": ["python-pptx", "zip_xml_parsing"],
                "image_formats_supported": ["PNG", "WMF", "EMF", "JPEG"],
                "table_extraction_enhanced": True,
                "table_position_tracking": True,
                "clip_descriptions_generated": True,
                "output_format": "JSON/CSV"
            }
        }

        # 保存元数据
        metadata_path = f"output/{base_filename}_pptx_metadata.json"
        with open(metadata_path, "w", encoding="utf-8") as f:
            json.dump(metadata, f, ensure_ascii=False, indent=2)

        return metadata_path


def process_pptx_file_advanced(preprocessor, file_path, fast_mode=False):
    """
    高级PPTX处理的入口函数

    Args:
        preprocessor: MultimodalPreprocessor实例
        file_path: PPTX文件路径
        fast_mode: 快速模式，跳过耗时处理
    """
    processor = AdvancedPPTProcessor(preprocessor, fast_mode=fast_mode)
    processor.process_pptx_file_advanced(file_path)

class MultimodalPreprocessor:
    def __init__(self):
        """初始化多模态预处理工具"""
        print("🚀 开始初始化多模态预处理工具...")

        # 检测设备
        print("📱 检测计算设备...")
        self.device = "cuda" if torch.cuda.is_available() else "cpu"
        print(f"✓ 使用设备: {self.device}")

        # 创建输出目录
        print("📁 创建输出目录...")
        os.makedirs("output/text", exist_ok=True)
        os.makedirs("output/images", exist_ok=True)
        os.makedirs("output/formulas", exist_ok=True)
        os.makedirs("output/tables", exist_ok=True)
        os.makedirs("output/code", exist_ok=True)
        print("✓ 输出目录创建完成")

        # 初始化CLIP模型（可能较慢）
        print("🤖 正在加载CLIP模型...")
        print("   ⏳ 本地模型加载中，请稍候...")
        # ===== 修改这里：使用你的本地CLIP模型路径 =====
        local_clip_path = r"F:\Models\clip-vit-base-patch32"
        try:
            if os.path.exists(local_clip_path):
                print(f"   📁 找到本地模型: {local_clip_path}")
                self.clip_model = CLIPModel.from_pretrained(local_clip_path, local_files_only=True).to(self.device)
                print("   ✓ 本地CLIP模型加载完成")
            else:
                print(f"   ❌ 本地模型路径不存在: {local_clip_path}")
                raise FileNotFoundError("本地模型不存在")
        except Exception as e:
            print(f"   ❌ 本地CLIP模型加载失败: {e}")
            print("   ⏳ 尝试在线下载CLIP模型，这可能需要几分钟...")
            try:
                self.clip_model = CLIPModel.from_pretrained("openai/clip-vit-base-patch32").to(self.device)
                print("   ✓ 在线CLIP模型下载并加载完成")
            except Exception as e2:
                print(f"   ❌ CLIP模型加载完全失败: {e2}")
                raise e2

        print("🔧 正在加载CLIP处理器...")
        print("   ⏳ 处理器加载中（可能需要下载）...")
        # ===== 修改这里：使用你的本地CLIP处理器路径 =====
        local_clip_path = r"F:\Models\clip-vit-base-patch32"
        try:
            if os.path.exists(local_clip_path):
                print(f"   📁 使用本地处理器: {local_clip_path}")
                self.clip_processor = CLIPProcessor.from_pretrained(local_clip_path, local_files_only=True)
                print("   ✓ 本地CLIP处理器加载完成")
            else:
                print(f"   ❌ 本地处理器路径不存在，使用在线版本")
                self.clip_processor = CLIPProcessor.from_pretrained("openai/clip-vit-base-patch32")
                print("   ✓ 在线CLIP处理器加载完成")
        except Exception as e:
            print(f"   ❌ CLIP处理器加载失败: {e}")
            raise e

        # 初始化OCR引擎（首次运行较慢）
        print("👁 正在初始化OCR引擎...")
        print("   ⏳ 首次运行需要下载模型文件，这可能需要几分钟，请耐心等待...")
        print("   📥 正在下载中文和英文OCR模型...")
        try:
            self.ocr_reader = easyocr.Reader(['ch_sim', 'en'], gpu=False)  # 强制使用CPU避免GPU问题
            print("   ✓ OCR引擎初始化完成")
        except Exception as e:
            print(f"   ⚠ OCR初始化失败，将跳过OCR功能: {e}")
            print("   将继续运行，但跳过OCR公式识别功能")
            self.ocr_reader = None

        # 存储处理结果
        self.results = []

        print("🎉 多模态预处理工具初始化完成！")
        print("=" * 50)

    def process_pdf(self, file_path):
        """处理PDF文件，提取文本和图像"""
        print(f"开始处理PDF文件: {file_path}")
        doc = fitz.open(file_path)
        base_filename = os.path.splitext(os.path.basename(file_path))[0]

        # 为当前PDF文件单独记录结果
        current_results = []
        original_results = self.results
        self.results = current_results

        try:
            for page_num in range(len(doc)):
                print(f"处理第 {page_num + 1}/{len(doc)} 页...")
                page = doc.load_page(page_num)
                page_text = page.get_text()

                # 处理页面图像
                image_list = page.get_images(full=True)
                page_images = []

                for img_index, img in enumerate(image_list):
                    xref = img[0]
                    base_image = doc.extract_image(xref)
                    image_bytes = base_image["image"]

                    # 保存原始图像
                    img_path = f"output/images/{base_filename}_p{page_num + 1}_img{img_index + 1}.{base_image['ext']}"
                    with open(img_path, "wb") as img_file:
                        img_file.write(image_bytes)

                    # 处理图像并保存
                    image_data = self.process_image(img_path, page_text)
                    self.save_image_data(image_data, base_filename, page_num, img_index)
                    page_images.append(img_path)

                # 处理页面文本
                text_data = self.process_text(page_text, page_num, page_images)
                text_data["source"] = f"{base_filename}_page{page_num + 1}"
                self.save_text_data(text_data, base_filename, page_num)

                # 提取页面中的公式、表格、代码
                self.extract_formulas_from_page(page, page_text, base_filename, page_num)
                self.extract_tables_from_page(page, page_text, base_filename, page_num)
                self.extract_code_from_page(page_text, base_filename, page_num)

            # 保存PDF专用元数据
            self.save_pdf_metadata(file_path, base_filename)
            print(f"PDF处理完成！结果保存在output/{base_filename}_pdf_metadata.json")

        finally:
            # 恢复原始结果列表并合并当前结果
            self.results = original_results
            self.results.extend(current_results)

    def process_text(self, text, page_num, page_images):
        """处理文本内容"""
        # 清理文本
        cleaned_text = text.strip()
        if not cleaned_text:
            cleaned_text = "[页面无文本内容]"

        # 使用CLIP生成文本语义向量
        try:
            inputs = self.clip_processor(text=cleaned_text, return_tensors="pt", padding=True, truncation=True)
            inputs = {k: v.to(self.device) for k, v in inputs.items()}

            with torch.no_grad():
                text_features = self.clip_model.get_text_features(**inputs)
                text_vector = text_features.cpu().numpy()[0]
        except Exception as e:
            print(f"文本向量化失败: {e}")
            text_vector = np.zeros(512)  # CLIP默认向量维度

        return {
            "type": "text",
            "page": page_num + 1,
            "raw_text": cleaned_text,
            "word_count": len(cleaned_text),
            "associated_images": page_images,
            "text_vector": text_vector.tolist()
        }

    def process_image(self, image_path, page_text):
        """处理图像（使用CLIP）"""
        # 图像增强
        enhanced_path = self.enhance_image(image_path)

        # 获取图像基本信息
        try:
            with Image.open(image_path) as img:
                width, height = img.size
                format_type = img.format
                mode = img.mode
        except Exception as e:
            print(f"读取图像信息失败: {e}")
            width = height = 0
            format_type = mode = "unknown"

        # 使用CLIP生成图像向量和描述
        try:
            image = Image.open(enhanced_path)
            inputs = self.clip_processor(images=image, return_tensors="pt").to(self.device)

            with torch.no_grad():
                image_features = self.clip_model.get_image_features(**inputs)
                image_vector = image_features.cpu().numpy()[0]

            # 生成图像描述标签
            description_tags = self.generate_image_descriptions(enhanced_path)

        except Exception as e:
            print(f"图像处理失败: {e}")
            image_vector = np.zeros(512)
            description_tags = []

        return {
            "type": "image",
            "image_path": image_path,
            "enhanced_path": enhanced_path,
            "width": width,
            "height": height,
            "format": format_type,
            "mode": mode,
            "page_text_context": page_text[:200] + "..." if len(page_text) > 200 else page_text,
            "image_vector": image_vector.tolist(),
            "clip_descriptions": description_tags
        }

    def enhance_image(self, image_path):
        """图像增强处理"""
        img = Image.open(image_path)

        # 对比度增强
        enhancer = ImageEnhance.Contrast(img)
        img = enhancer.enhance(1.5)

        # 锐度增强
        enhancer = ImageEnhance.Sharpness(img)
        img = enhancer.enhance(2.0)

        # 保存增强后的图像
        enhanced_path = image_path.replace(".", "_enhanced.")
        img.save(enhanced_path)

        return enhanced_path

    def clip_generate_description(self, image_path: str) -> str:
        """基于CLIP为图片生成描述文本并保存为JSON，返回描述文件路径。"""
        try:
            descriptions = self.generate_image_descriptions(image_path)
        except Exception as e:
            print(f"生成图片描述失败: {e}")
            descriptions = []

        base, _ = os.path.splitext(os.path.basename(image_path))
        desc_path = os.path.join("output", "images", f"{base}_desc.json")
        try:
            with open(desc_path, "w", encoding="utf-8") as f:
                json.dump({
                    "image_path": image_path,
                    "clip_descriptions": descriptions
                }, f, ensure_ascii=False, indent=2)
        except Exception as e:
            print(f"保存图片描述失败: {e}")
        return desc_path

    def generate_image_descriptions(self, image_path):
        """使用CLIP生成图像描述标签"""
        try:
            image = Image.open(image_path)
            image_input = self.clip_processor(images=image, return_tensors="pt").to(self.device)

            # 预定义的描述候选列表
            text_descriptions = [
                "科学图表", "数学公式", "数据图表", "流程图",
                "实验装置", "分子结构", "几何图形", "统计图表",
                "技术示意图", "概念图", "网络图", "系统架构图",
                "照片", "插图", "示例图", "对比图",
                "文本图像", "表格", "代码", "截图"
            ]

            text_inputs = self.clip_processor(
                text=text_descriptions,
                return_tensors="pt",
                padding=True,
                truncation=True
            ).to(self.device)

            with torch.no_grad():
                image_features = self.clip_model.get_image_features(**image_input)
                text_features = self.clip_model.get_text_features(**text_inputs)

                # 归一化特征向量
                image_features = image_features / image_features.norm(dim=-1, keepdim=True)
                text_features = text_features / text_features.norm(dim=-1, keepdim=True)

                # 计算相似度
                similarity = (image_features @ text_features.T).softmax(dim=-1)
                values, indices = similarity[0].topk(5)  # 取前5个最相似的描述

                descriptions = []
                for value, idx in zip(values.cpu().numpy(), indices.cpu().numpy()):
                    descriptions.append({
                        "description": text_descriptions[idx],
                        "confidence": float(value)
                    })

            return descriptions

        except Exception as e:
            print(f"图像描述生成错误: {image_path}, {str(e)}")
            return []

    def save_text_data(self, data, filename, page_num):
        """保存文本处理结果"""
        output_path = f"output/text/{filename}_p{page_num + 1}.json"
        with open(output_path, "w", encoding="utf-8") as f:
            json.dump(data, f, ensure_ascii=False, indent=2)

        self.results.append({
            "type": "text",
            "page": page_num + 1,
            "file": output_path
        })

    def save_image_data(self, data, filename, page_num, img_index):
        """保存图像处理结果"""
        output_path = f"output/images/{filename}_p{page_num + 1}_img{img_index + 1}.json"
        with open(output_path, "w", encoding="utf-8") as f:
            json.dump(data, f, ensure_ascii=False, indent=2)

        self.results.append({
            "type": "image",
            "page": page_num + 1,
            "file": output_path
        })

    def save_pdf_metadata(self, file_path, filename):
        """保存PDF专用元数据文件"""
        # 统计信息
        text_files = [r for r in self.results if r["type"] == "text"]
        image_files = [r for r in self.results if r["type"] == "image"]
        formula_files = [r for r in self.results if r["type"] == "formula"]
        table_files = [r for r in self.results if r["type"] == "table"]
        code_files = [r for r in self.results if r["type"] == "code"]

        # 计算表格统计
        total_table_rows = sum(r.get("rows", 0) for r in table_files)
        total_table_columns = sum(r.get("columns", 0) for r in table_files)

        # 计算代码统计
        total_code_lines = sum(r.get("line_count", 0) for r in code_files)
        code_languages = list(set(r.get("language", "unknown") for r in code_files))

        metadata = {
            "processing_date": datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
            "source_file": filename,
            "source_path": file_path,
            "file_type": "PDF",
            "processing_method": "pymupdf_ocr_camelot",
            "statistics": {
                "total_pages": len(set(r["page"] for r in self.results)),
                "total_text_blocks": len(text_files),
                "total_images": len(image_files),
                "total_formulas": len(formula_files),
                "total_tables": len(table_files),
                "total_table_rows": total_table_rows,
                "total_table_columns": total_table_columns,
                "total_code_blocks": len(code_files),
                "total_code_lines": total_code_lines,
                "code_languages": code_languages
            },
            "files": {
                "text_files": text_files,
                "image_files": image_files,
                "formula_files": formula_files,
                "table_files": table_files,
                "code_files": code_files
            },
            "processing_info": {
                "clip_model": "openai/clip-vit-base-patch32",
                "device": self.device,
                "output_format": "JSON/CSV",
                "ocr_enabled": self.ocr_reader is not None,
                "extraction_features": ["text", "images", "formulas", "tables", "code"]
            }
        }

        with open(f"output/{filename}_pdf_metadata.json", "w", encoding="utf-8") as f:
            json.dump(metadata, f, ensure_ascii=False, indent=2)

    def extract_formulas_from_page(self, page, page_text, filename, page_num):
        """从页面中提取数学公式"""
        formulas = []

        # 1. 从文本中提取LaTeX格式的公式
        latex_patterns = [
            r'\$\$([^$]+)\$\$',  # 块级公式 $$...$$
            r'\$([^$]+)\$',  # 行内公式 $...$
            r'\\begin\{equation\}(.*?)\\end\{equation\}',  # equation环境
            r'\\begin\{align\}(.*?)\\end\{align\}',  # align环境
            r'\\begin\{math\}(.*?)\\end\{math\}',  # math环境
        ]

        for i, pattern in enumerate(latex_patterns):
            matches = re.findall(pattern, page_text, re.DOTALL)
            for j, match in enumerate(matches):
                formula_data = {
                    "type": "formula",
                    "page": page_num + 1,
                    "formula_id": f"{filename}_p{page_num + 1}_formula{len(formulas) + 1}",
                    "content": match.strip(),
                    "format": "latex",
                    "extraction_method": f"pattern_{i + 1}",
                    "context": self.get_text_context(page_text, match, 100)
                }
                formulas.append(formula_data)

        # 2. 从图像中识别公式（使用OCR）- 简化版本避免卡住
        if self.ocr_reader and len(formulas) < 5:  # 限制OCR处理，避免卡住
            try:
                # 获取页面图像
                pix = page.get_pixmap()
                img_data = pix.tobytes("png")
                img_array = np.frombuffer(img_data, np.uint8)
                img = cv2.imdecode(img_array, cv2.IMREAD_COLOR)

                # OCR识别，设置超时
                results = self.ocr_reader.readtext(img, width_ths=0.7, height_ths=0.7)

                for result in results[:3]:  # 只处理前3个结果，避免过多处理
                    text = result[1]
                    confidence = result[2]

                    # 检查是否包含数学符号
                    math_symbols = ['∑', '∫', '∂', '∆', '∇', '∞', '±', '≠', '≤', '≥', 'α', 'β', 'γ', 'δ', 'θ', 'λ', 'μ',
                                    'π', 'σ', 'φ', 'ψ', 'ω']
                    if any(symbol in text for symbol in math_symbols) and confidence > 0.6:
                        if any(char.isdigit() or char in '+-*/=()[]{}^_' for char in text):
                            formula_data = {
                                "type": "formula",
                                "page": page_num + 1,
                                "formula_id": f"{filename}_p{page_num + 1}_formula{len(formulas) + 1}",
                                "content": text,
                                "format": "ocr_text",
                                "confidence": float(confidence),
                                "extraction_method": "ocr",
                                "bbox": [[float(pt[0]), float(pt[1])] for pt in result[0]]  # 转换为Python原生类型
                            }
                            formulas.append(formula_data)

            except Exception as e:
                print(f"OCR公式识别失败 (页面 {page_num + 1}): {e}")

        # 保存公式数据
        if formulas:
            self.save_formulas_data(formulas, filename, page_num)

        return formulas

    def extract_tables_from_page(self, page, page_text, filename, page_num):
        """从页面中提取表格"""
        tables = []

        try:
            # 使用camelot提取表格 - 限制处理时间，避免卡住
            pdf_path = None
            for file in os.listdir("input"):
                if file.lower().endswith('.pdf') and filename in file:
                    pdf_path = os.path.join("input", file)
                    break

            if pdf_path and os.path.exists(pdf_path) and page_num < 10:  # 只处理前10页，避免卡住
                # 提取当前页面的表格，限制处理
                camelot_tables = camelot.read_pdf(pdf_path, pages=str(page_num + 1), flavor='lattice')

                for i, table in enumerate(camelot_tables[:2]):  # 只处理前2个表格
                    if table.df is not None and not table.df.empty and len(table.df) > 0:
                        table_data = {
                            "type": "table",
                            "page": page_num + 1,
                            "table_id": f"{filename}_p{page_num + 1}_table{i + 1}",
                            "rows": len(table.df),
                            "columns": len(table.df.columns),
                            "data": table.df.to_dict('records'),
                            "extraction_method": "camelot",
                            "accuracy": getattr(table, 'accuracy', 0.0)
                        }
                        tables.append(table_data)

        except Exception as e:
            print(f"Camelot表格提取跳过 (页面 {page_num + 1}): {e}")

        # 备用方法：从文本中识别表格模式
        table_patterns = self.detect_text_tables(page_text)
        for i, pattern in enumerate(table_patterns):
            table_data = {
                "type": "table",
                "page": page_num + 1,
                "table_id": f"{filename}_p{page_num + 1}_texttable{i + 1}",
                "content": pattern,
                "extraction_method": "text_pattern",
                "context": self.get_text_context(page_text, pattern, 50)
            }
            tables.append(table_data)

        # 保存表格数据
        if tables:
            self.save_tables_data(tables, filename, page_num)

        return tables

    def extract_code_from_page(self, page_text, filename, page_num):
        """从页面文本中提取代码块"""
        code_blocks = []

        # 代码块模式
        code_patterns = [
            r'```(\w*)\n(.*?)```',  # Markdown代码块
            r'`([^`]+)`',  # 行内代码
            r'(?:^|\n)((?:    |\t)[^\n]+(?:\n(?:    |\t)[^\n]+)*)',  # 缩进代码块
        ]

        # 编程语言关键字
        programming_keywords = [
            'def ', 'class ', 'import ', 'from ', 'if ', 'else:', 'for ', 'while ', 'return',
            'function', 'var ', 'let ', 'const ', 'console.log', 'print(', 'println',
            'public ', 'private ', 'static ', 'void ', 'int ', 'String ', 'boolean',
            '#include', 'using namespace', 'int main(', 'printf(', 'cout <<'
        ]

        for i, pattern in enumerate(code_patterns):
            matches = re.findall(pattern, page_text, re.DOTALL | re.MULTILINE)

            for j, match in enumerate(matches):
                if isinstance(match, tuple):
                    language = match[0] if match[0] else "unknown"
                    content = match[1] if len(match) > 1 else match[0]
                else:
                    content = match
                    language = "unknown"

                # 检查是否包含编程关键字
                if any(keyword in content for keyword in programming_keywords) or len(content.strip()) > 20:
                    code_data = {
                        "type": "code",
                        "page": page_num + 1,
                        "code_id": f"{filename}_p{page_num + 1}_code{len(code_blocks) + 1}",
                        "content": content.strip(),
                        "language": language,
                        "extraction_method": f"pattern_{i + 1}",
                        "line_count": len(content.strip().split('\n')),
                        "context": self.get_text_context(page_text, content, 100)
                    }
                    code_blocks.append(code_data)

        # 保存代码数据
        if code_blocks:
            self.save_code_data(code_blocks, filename, page_num)

        return code_blocks

    def get_text_context(self, full_text, target_text, context_length=100):
        """获取目标文本的上下文"""
        try:
            index = full_text.find(target_text)
            if index == -1:
                return target_text

            start = max(0, index - context_length)
            end = min(len(full_text), index + len(target_text) + context_length)
            return full_text[start:end]
        except:
            return target_text

    def detect_text_tables(self, text):
        """从文本中检测表格模式"""
        tables = []
        lines = text.split('\n')

        # 寻找包含多个制表符或空格分隔的行
        table_lines = []
        for line in lines:
            # 检查是否包含表格特征：多个制表符、竖线分隔符等
            if '\t' in line and line.count('\t') >= 2:
                table_lines.append(line)
            elif '|' in line and line.count('|') >= 2:
                table_lines.append(line)
            elif re.search(r'\s{3,}', line) and len(line.split()) >= 3:
                table_lines.append(line)
            else:
                if table_lines and len(table_lines) >= 2:
                    tables.append('\n'.join(table_lines))
                table_lines = []

        # 检查最后一组
        if table_lines and len(table_lines) >= 2:
            tables.append('\n'.join(table_lines))

        return tables

    def save_formulas_data(self, formulas, filename, page_num):
        """保存公式数据"""
        # JSON格式保存
        json_path = f"output/formulas/{filename}_p{page_num + 1}_formulas.json"
        with open(json_path, "w", encoding="utf-8") as f:
            json.dump(formulas, f, ensure_ascii=False, indent=2)

        # CSV格式保存
        csv_path = f"output/formulas/{filename}_p{page_num + 1}_formulas.csv"
        if formulas:
            df = pd.DataFrame(formulas)
            df.to_csv(csv_path, index=False, encoding="utf-8")

        # 记录到结果
        for formula in formulas:
            self.results.append({
                "type": "formula",
                "page": page_num + 1,
                "file": json_path,
                "formula_id": formula["formula_id"]
            })

    def save_tables_data(self, tables, filename, page_num):
        """保存表格数据"""
        # JSON格式保存
        json_path = f"output/tables/{filename}_p{page_num + 1}_tables.json"
        with open(json_path, "w", encoding="utf-8") as f:
            json.dump(tables, f, ensure_ascii=False, indent=2)

        # 为每个表格单独保存CSV
        for i, table in enumerate(tables):
            if table.get("data") and isinstance(table["data"], list):
                csv_path = f"output/tables/{table['table_id']}.csv"
                try:
                    df = pd.DataFrame(table["data"])
                    df.to_csv(csv_path, index=False, encoding="utf-8")
                except Exception as e:
                    print(f"保存表格CSV失败: {e}")

        # 记录到结果
        for table in tables:
            self.results.append({
                "type": "table",
                "page": page_num + 1,
                "file": json_path,
                "table_id": table["table_id"],
                "rows": table.get("rows", 0),
                "columns": table.get("columns", 0)
            })

    def save_code_data(self, code_blocks, filename, page_num):
        """保存代码数据"""
        # JSON格式保存
        json_path = f"output/code/{filename}_p{page_num + 1}_code.json"
        with open(json_path, "w", encoding="utf-8") as f:
            json.dump(code_blocks, f, ensure_ascii=False, indent=2)

        # CSV格式保存
        csv_path = f"output/code/{filename}_p{page_num + 1}_code.csv"
        if code_blocks:
            df = pd.DataFrame(code_blocks)
            df.to_csv(csv_path, index=False, encoding="utf-8")

        # 为每个代码块单独保存文件
        for code in code_blocks:
            if code.get("language") and code.get("language") != "unknown":
                ext = self.get_file_extension(code["language"])
                code_file_path = f"output/code/{code['code_id']}.{ext}"
                with open(code_file_path, "w", encoding="utf-8") as f:
                    f.write(code["content"])

        # 记录到结果
        for code in code_blocks:
            self.results.append({
                "type": "code",
                "page": page_num + 1,
                "file": json_path,
                "code_id": code["code_id"],
                "language": code.get("language", "unknown"),
                "line_count": code.get("line_count", 0)
            })

    def get_file_extension(self, language):
        """根据编程语言获取文件扩展名"""
        extensions = {
            "python": "py",
            "javascript": "js",
            "java": "java",
            "cpp": "cpp",
            "c": "c",
            "csharp": "cs",
            "php": "php",
            "ruby": "rb",
            "go": "go",
            "rust": "rs",
            "swift": "swift",
            "kotlin": "kt",
            "typescript": "ts",
            "html": "html",
            "css": "css",
            "sql": "sql",
            "shell": "sh",
            "bash": "sh",
            "powershell": "ps1"
        }
        return extensions.get(language.lower(), "txt")


def process_multimodal_files(input_dir, output_dir="output", clip_model_path=None, fast_mode=False, file_types=None):
    """
    多模态文件处理主函数

    Args:
        input_dir (str): 输入文件目录路径
        output_dir (str): 输出目录路径，默认为"output"
        clip_model_path (str): 本地CLIP模型路径，如果为None则使用在线模型
        fast_mode (bool): 快速模式，跳过耗时的CLIP描述生成
        file_types (list): 支持的文件类型列表，默认为['pdf', 'pptx']

    Returns:
        dict: 处理结果统计信息
    """
    import sys
    import os

    # 获取当前文件的绝对路径
    current_dir = os.path.dirname(os.path.abspath(__file__))
    parent_dir = os.path.dirname(current_dir)

    # 添加路径
    sys.path.append(parent_dir)  # 项目根目录
    sys.path.append(current_dir)  # src目录

    import json
    import fitz  # PyMuPDF
    import torch
    import numpy as np
    from PIL import Image, ImageEnhance
    from transformers import CLIPProcessor, CLIPModel
    import datetime
    import pandas as pd
    import cv2
    import easyocr
    import re
    import pdfplumber
    import camelot
    import csv

    # 设置默认文件类型
    if file_types is None:
        file_types = ['pdf', 'pptx']

    print("=" * 60)
    print("多模态文件数据预处理器")
    print("=" * 60)
    print(f"输入目录: {input_dir}")
    print(f"输出目录: {output_dir}")
    print(f"快速模式: {fast_mode}")
    print(f"支持文件类型: {file_types}")

    # 确保输出目录存在
    os.makedirs(f"{output_dir}/text", exist_ok=True)
    os.makedirs(f"{output_dir}/images", exist_ok=True)
    os.makedirs(f"{output_dir}/formulas", exist_ok=True)
    os.makedirs(f"{output_dir}/tables", exist_ok=True)
    os.makedirs(f"{output_dir}/code", exist_ok=True)

    try:
        # 初始化处理器
        print("⚙️ 正在初始化处理器...")
        print("⚠️ 注意：首次运行可能需要下载模型，请耐心等待...")

        # 创建临时的MultimodalPreprocessor类实例
        class TempMultimodalPreprocessor:
            def __init__(self, clip_model_path=None, output_dir="output"):
                self.output_dir = output_dir
                self.device = "cuda" if torch.cuda.is_available() else "cpu"
                print(f"✓ 使用设备: {self.device}")

                # 初始化CLIP模型
                print("🤖 正在加载CLIP模型...")
                if clip_model_path and os.path.exists(clip_model_path):
                    print(f"   📁 使用本地模型: {clip_model_path}")
                    self.clip_model = CLIPModel.from_pretrained(clip_model_path, local_files_only=True).to(self.device)
                    self.clip_processor = CLIPProcessor.from_pretrained(clip_model_path, local_files_only=True)
                    print("   ✓ 本地CLIP模型加载完成")
                else:
                    print("   ⏳ 使用在线模型...")
                    self.clip_model = CLIPModel.from_pretrained("openai/clip-vit-base-patch32").to(self.device)
                    self.clip_processor = CLIPProcessor.from_pretrained("openai/clip-vit-base-patch32")
                    print("   ✓ 在线CLIP模型加载完成")

                # 初始化OCR引擎
                print("👁 正在初始化OCR引擎...")
                try:
                    self.ocr_reader = easyocr.Reader(['ch_sim', 'en'], gpu=False)
                    print("   ✓ OCR引擎初始化完成")
                except Exception as e:
                    print(f"   ⚠ OCR初始化失败，将跳过OCR功能: {e}")
                    self.ocr_reader = None

                self.results = []

        # 使用原有的MultimodalPreprocessor类，但修改输出路径
        processor = MultimodalPreprocessor()
        # 更新输出路径
        processor.output_text_dir = f"{output_dir}/text"
        processor.output_table_dir = f"{output_dir}/tables"
        processor.output_img_dir = f"{output_dir}/images"

        print("✅ 处理器初始化完成!")

        # 检查输入目录
        if not os.path.exists(input_dir):
            raise FileNotFoundError(f"输入目录不存在: {input_dir}")

        # 查找支持的文件
        file_extensions = tuple(f'.{ft}' for ft in file_types)
        input_files = [f for f in os.listdir(input_dir)
                       if f.lower().endswith(file_extensions) and not f.startswith('~$')]

        if not input_files:
            print(f"❌ 在 {input_dir} 目录中未找到支持的文件")
            print(f"支持的文件类型: {file_types}")
            return {
                "status": "no_files",
                "processed_files": 0,
                "total_files": 0,
                "file_types": file_types
            }

        # 统计文件类型
        file_stats = {}
        for ft in file_types:
            count = sum(1 for f in input_files if f.lower().endswith(f'.{ft}'))
            file_stats[ft] = count
            if count > 0:
                print(f"✓ 找到 {count} 个 {ft.upper()} 文件")

        # 处理文件
        processed_count = 0
        failed_files = []

        for idx, in_file in enumerate(input_files, 1):
            input_path = os.path.join(input_dir, in_file)
            print(f"\n📄 [{idx}/{len(input_files)}] 开始处理: {in_file}")

            try:
                if in_file.lower().endswith('.pdf'):
                    print("   📚 使用PDF处理器（PyMuPDF + OCR + 表格提取）...")
                    processor.process_pdf(input_path)
                    processed_count += 1

                elif in_file.lower().endswith('.pptx'):
                    print("   📊 使用高级PPTX处理器（ZIP+XML解析 + 优化表格提取）...")
                    try:
                        # 使用高级PPTX处理器
                        advanced_processor = AdvancedPPTProcessor(processor, fast_mode=fast_mode)
                        advanced_processor.process_pptx_file_advanced(input_path)
                        processed_count += 1
                    except Exception as e:
                        print(f"   ⚠ 高级PPTX处理失败，尝试基础处理: {e}")
                        # 这里可以添加基础PPTX处理逻辑
                        failed_files.append({"file": in_file, "error": str(e)})

                print(f"   ✅ [{idx}/{len(input_files)}] 完成处理: {in_file}")

            except Exception as e:
                print(f"   ❌ [{idx}/{len(input_files)}] 处理失败: {in_file}")
                print(f"   错误: {e}")
                failed_files.append({"file": in_file, "error": str(e)})

        # 生成处理报告
        processing_report = {
            "status": "completed",
            "processing_date": datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
            "input_directory": input_dir,
            "output_directory": output_dir,
            "total_files": len(input_files),
            "processed_files": processed_count,
            "failed_files": len(failed_files),
            "file_statistics": file_stats,
            "failed_file_details": failed_files,
            "settings": {
                "fast_mode": fast_mode,
                "clip_model_path": clip_model_path,
                "supported_file_types": file_types,
                "device": processor.device
            },
            "results_summary": {
                "total_results": len(processor.results),
                "text_extractions": len([r for r in processor.results if r["type"] in ["text", "ppt_text"]]),
                "image_extractions": len(
                    [r for r in processor.results if r["type"] in ["image", "ppt_image", "ppt_image_zip"]]),
                "table_extractions": len([r for r in processor.results if r["type"] in ["table", "ppt_table"]]),
                "formula_extractions": len([r for r in processor.results if r["type"] == "formula"]),
                "code_extractions": len([r for r in processor.results if r["type"] == "code"])
            }
        }

        # 保存处理报告
        report_path = os.path.join(output_dir, "processing_report.json")
        with open(report_path, "w", encoding="utf-8") as f:
            json.dump(processing_report, f, ensure_ascii=False, indent=2)

        print(f"\n🎉 处理完成！")
        print(f"📊 成功处理: {processed_count}/{len(input_files)} 个文件")
        if failed_files:
            print(f"⚠️ 失败文件: {len(failed_files)} 个")
        print(f"📁 结果保存在: {output_dir}")
        print(f"📋 处理报告: {report_path}")

        return processing_report

    except Exception as e:
        print(f"❌ 处理过程中发生严重错误: {e}")
        import traceback
        traceback.print_exc()

        return {
            "status": "error",
            "error": str(e),
            "processed_files": 0,
            "total_files": 0
        }


# 使用示例
if __name__ == "__main__":
    # 示例调用
    result = process_multimodal_files(
        input_dir=r"C:\Users\Lin\PycharmProjects\PythonProject\input",
        output_dir="output",
        clip_model_path=r"F:\Models\clip-vit-base-patch32",  # 可选：本地CLIP模型路径
        fast_mode=False,  # 设置为True可跳过耗时的CLIP描述生成
        file_types=['pdf', 'pptx']  # 支持的文件类型
    )

    print(f"\n最终处理结果: {result['status']}")
    print(f"处理文件数: {result['processed_files']}/{result['total_files']}")
